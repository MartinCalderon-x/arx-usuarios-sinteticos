"""Reports API routes."""
from fastapi import APIRouter, HTTPException, Depends
from fastapi.responses import FileResponse, StreamingResponse
from typing import Optional
from pydantic import BaseModel
from datetime import datetime
import io

from app.core.security import require_auth
from app.services import supabase

router = APIRouter()


class ReporteRequest(BaseModel):
    """Schema for report generation request."""
    titulo: str
    formato: str = "pdf"  # pdf or pptx
    arquetipos_ids: Optional[list[str]] = None
    analisis_ids: Optional[list[str]] = None
    sesiones_ids: Optional[list[str]] = None
    incluir_resumen: bool = True
    incluir_recomendaciones: bool = True


class ReporteResponse(BaseModel):
    """Schema for report response."""
    id: str
    titulo: str
    formato: str
    url: Optional[str] = None
    created_at: str


def get_user_id(user: dict) -> str:
    """Extract user ID from JWT payload."""
    return user.get("sub")


@router.post("/generar", response_model=ReporteResponse)
async def generar_reporte(request: ReporteRequest, user: dict = Depends(require_auth)):
    """Generate a report from analysis and interaction data. Requires authentication."""
    user_id = get_user_id(user)

    # Collect data for report
    contenido = {
        "arquetipos": [],
        "analisis": [],
        "sesiones": []
    }

    # Get archetipos data
    if request.arquetipos_ids:
        for arq_id in request.arquetipos_ids:
            arq = await supabase.get_arquetipo(arq_id, user_id)
            if arq:
                contenido["arquetipos"].append(arq)

    # Get analysis data
    if request.analisis_ids:
        for ana_id in request.analisis_ids:
            ana = await supabase.get_analisis(ana_id, user_id)
            if ana:
                contenido["analisis"].append(ana)

    # Get session data with messages
    if request.sesiones_ids:
        for ses_id in request.sesiones_ids:
            ses = await supabase.get_sesion(ses_id, user_id)
            if ses:
                historial = await supabase.get_historial(ses_id)
                ses["mensajes"] = historial
                contenido["sesiones"].append(ses)

    # Create report record
    data = {
        "titulo": request.titulo,
        "formato": request.formato,
        "contenido": contenido,
        "arquetipos_ids": request.arquetipos_ids or [],
        "analisis_ids": request.analisis_ids or [],
        "sesiones_ids": request.sesiones_ids or []
    }

    result = await supabase.create_reporte(data, user_id)
    if not result:
        raise HTTPException(status_code=500, detail="Error al crear reporte")

    return ReporteResponse(
        id=result["id"],
        titulo=result["titulo"],
        formato=result["formato"],
        url=f"/api/reportes/{result['id']}/descargar",
        created_at=result["created_at"]
    )


@router.get("/{reporte_id}/descargar")
async def descargar_reporte(reporte_id: str, user: dict = Depends(require_auth)):
    """Download a generated report. Requires authentication."""
    user_id = get_user_id(user)

    # Get report
    reporte = await supabase.get_reporte(reporte_id, user_id)
    if not reporte:
        raise HTTPException(status_code=404, detail="Reporte no encontrado")

    formato = reporte.get("formato", "pdf")
    contenido = reporte.get("contenido", {})
    titulo = reporte.get("titulo", "Reporte")

    if formato == "pdf":
        # Generate PDF using ReportLab
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.styles import getSampleStyleSheet

        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        # Title
        story.append(Paragraph(titulo, styles['Title']))
        story.append(Spacer(1, 20))

        # Archetipos section
        if contenido.get("arquetipos"):
            story.append(Paragraph("Arquetipos", styles['Heading1']))
            for arq in contenido["arquetipos"]:
                story.append(Paragraph(f"<b>{arq.get('nombre', 'N/A')}</b>", styles['Heading2']))
                story.append(Paragraph(f"Descripcion: {arq.get('descripcion', 'N/A')}", styles['Normal']))
                story.append(Paragraph(f"Edad: {arq.get('edad', 'N/A')} | Ocupacion: {arq.get('ocupacion', 'N/A')}", styles['Normal']))
                story.append(Spacer(1, 10))

        # Analysis section
        if contenido.get("analisis"):
            story.append(Paragraph("Analisis Visual", styles['Heading1']))
            for ana in contenido["analisis"]:
                story.append(Paragraph(f"Clarity Score: {ana.get('clarity_score', 'N/A')}", styles['Normal']))
                insights = ana.get("insights", [])
                if insights:
                    story.append(Paragraph("Insights:", styles['Heading3']))
                    for insight in insights:
                        story.append(Paragraph(f"- {insight}", styles['Normal']))
                story.append(Spacer(1, 10))

        # Sessions section
        if contenido.get("sesiones"):
            story.append(Paragraph("Sesiones de Interaccion", styles['Heading1']))
            for ses in contenido["sesiones"]:
                story.append(Paragraph(f"Sesion: {ses.get('id', 'N/A')[:8]}...", styles['Heading2']))
                mensajes = ses.get("mensajes", [])
                for msg in mensajes[:10]:  # Limit to 10 messages
                    rol = "Usuario" if msg.get("rol") == "usuario" else "Sintetico"
                    story.append(Paragraph(f"<b>{rol}:</b> {msg.get('contenido', '')[:200]}", styles['Normal']))
                story.append(Spacer(1, 10))

        doc.build(story)
        buffer.seek(0)

        return StreamingResponse(
            buffer,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={titulo}.pdf"}
        )

    elif formato == "pptx":
        # Generate PPTX using python-pptx
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()

        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = titulo
        subtitle.text = f"Generado: {datetime.now().strftime('%Y-%m-%d %H:%M')}"

        # Archetipos slides
        if contenido.get("arquetipos"):
            for arq in contenido["arquetipos"]:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                body = slide.placeholders[1]
                title.text = f"Arquetipo: {arq.get('nombre', 'N/A')}"
                tf = body.text_frame
                tf.text = arq.get('descripcion', '')
                p = tf.add_paragraph()
                p.text = f"Edad: {arq.get('edad', 'N/A')} | Ocupacion: {arq.get('ocupacion', 'N/A')}"

        # Analysis slides
        if contenido.get("analisis"):
            for ana in contenido["analisis"]:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                body = slide.placeholders[1]
                title.text = "Analisis Visual"
                tf = body.text_frame
                tf.text = f"Clarity Score: {ana.get('clarity_score', 'N/A')}"
                for insight in ana.get("insights", [])[:5]:
                    p = tf.add_paragraph()
                    p.text = f"- {insight}"

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        return StreamingResponse(
            buffer,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={titulo}.pptx"}
        )

    raise HTTPException(status_code=400, detail="Formato no soportado")


@router.get("/")
async def list_reportes(user: dict = Depends(require_auth)):
    """List generated reports. Requires authentication."""
    user_id = get_user_id(user)
    reportes, total = await supabase.list_reportes(user_id)
    return {"reportes": reportes, "total": total}


@router.delete("/{reporte_id}")
async def delete_reporte(reporte_id: str, user: dict = Depends(require_auth)):
    """Delete a report. Requires authentication."""
    user_id = get_user_id(user)
    deleted = await supabase.delete_reporte(reporte_id, user_id)
    if not deleted:
        raise HTTPException(status_code=404, detail="Reporte no encontrado")
    return {"message": "Reporte eliminado"}
